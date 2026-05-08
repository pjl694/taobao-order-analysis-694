from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(7.5)

def add_title(slide, text, size=24, bold=True, color=RGBColor(0, 51, 102)):
    title = slide.shapes.title
    title.text = text
    title.text_frame.paragraphs[0].font.size = Pt(size)
    title.text_frame.paragraphs[0].font.bold = bold
    title.text_frame.paragraphs[0].font.color.rgb = color

def add_text_box(slide, text, left, top, width, height, size=16, bold=False, color=RGBColor(0,0,0), bg_color=None):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.text = text
    tf.paragraphs[0].font.size = Pt(size)
    tf.paragraphs[0].font.bold = bold
    tf.paragraphs[0].font.color.rgb = color
    if bg_color:
        box.fill.solid()
        box.fill.fore_color.rgb = bg_color
    return box

def add_table(slide, rows, cols, data, left, top, width, height):
    table = slide.shapes.add_table(rows, cols, Inches(left), Inches(top), Inches(width), Inches(height)).table
    for i, row in enumerate(data):
        for j, cell in enumerate(row):
            table.cell(i, j).text = str(cell)
            table.cell(i, j).text_frame.paragraphs[0].font.size = Pt(14)
    return table

# 第1页
slide = prs.slides.add_slide(prs.slide_layouts[0])
add_title(slide, "运营商用户行为分析与用户分群", size=32)
add_text_box(slide, "基于SQL + Python + SPSS + Power BI的数据分析项目", 1, 3, 8, 1, size=18, color=RGBColor(100,100,100))
add_text_box(slide, "蒲俊霖 · 数据科学与大数据技术 · 2026年5月", 2.5, 6, 5, 0.5, size=16, color=RGBColor(100,100,100))

# 第2页
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide, "项目背景")
content = ("本项目模拟运营商用户行为分析场景，基于某电商平台真实用户行为数据\n\n"
           "通过对用户活跃行为的深入分析，构建用户价值分层体系\n\n"
           "目标是识别高价值用户、发现活跃规律、输出可落地的运营策略\n\n"
           "数据来源：天池淘宝用户行为数据集（12,256,906条记录）")
add_text_box(slide, content, 0.5, 1.5, 9, 5, size=18)

# 第3页
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide, "数据概览")
data = [["维度", "详情"],
        ["原始数据量", "12,256,906 条"],
        ["去重后", "8,164,040 条"],
        ["分析样本", "1%随机抽样（81,640条）"],
        ["时间范围", "2014年11月18日 - 12月18日"],
        ["核心字段", "user_id、behavior_type、time、item_category"]]
add_table(slide, 6, 2, data, 1, 1.5, 8, 2.5)
add_text_box(slide, "技术栈：Python(Pandas) / SPSS / Power BI", 6.5, 5.5, 3, 0.8, size=12, bg_color=RGBColor(230,230,230))

# 第4页
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide, "用户价值分层结果")
add_text_box(slide, "左栏：饼图占位符（插入Power BI饼图截图）", 0.5, 1.5, 4, 1.5, size=14)
add_text_box(slide, "右栏：柱状图占位符（插入Power BI柱状图截图）", 5.5, 1.5, 4, 1.5, size=14)
data2 = [["等级", "人数", "占比", "平均活跃天数"],
         ["高价值用户", "3,034", "34.4%", "10天"],
         ["中价值用户", "3,271", "37.1%", "5天"],
         ["低价值用户", "2,508", "28.5%", "2天以下"]]
add_table(slide, 4, 4, data2, 1, 4, 8, 1.5)

# 第5页
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide, "统计显著性验证 - SPSS分析")
add_text_box(slide, "左栏：ANOVA表截图占位符", 0.5, 1.5, 4, 2, size=14)
add_text_box(slide, "F值：11,229.94\np值：<0.001\nEta方：0.718", 5.5, 1.5, 4, 2, size=16, bold=True)
add_text_box(slide, "不同用户等级的活跃天数存在极显著差异，证明用户价值分层在统计学上是合理且有效的。", 0.5, 5, 9, 1, size=16, bold=True, bg_color=RGBColor(0,51,102), color=RGBColor(255,255,255))

# 第6页
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide, "活跃趋势与时段偏好")
add_text_box(slide, "DAU趋势折线图占位符", 0.5, 1.5, 4, 2, size=14)
add_text_box(slide, "时段活跃柱状图占位符", 5.5, 1.5, 4, 2, size=14)
add_text_box(slide, "晚间20-23时为用户活跃最高峰，运营活动建议集中在此时段，可最大化触达率。", 0.5, 5, 9, 1, size=16, bold=True, bg_color=RGBColor(230,230,230))

# 第7页
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide, "核心发现与运营建议")
data3 = [["发现", "运营建议"],
         ["高价值用户占比34.4%，是核心资产", "设计VIP权益、专属客服、优先体验"],
         ["中价值用户占比37.1%，有成长空间", "通过签到激励、任务引导提升活跃频次"],
         ["低价值用户占比28.5%，存在流失风险", "分析流失原因，推送召回优惠券"],
         ["晚间20-23时为活跃高峰", "营销推送、直播活动优先排布此时段"],
         ["用户分层经SPSS验证显著有效", "可据此设计差异化运营策略"]]
add_table(slide, 6, 2, data3, 0.5, 1.5, 9, 3)

# 第8页
slide = prs.slides.add_slide(prs.slide_layouts[1])
add_title(slide, "感谢阅读", size=44)
add_text_box(slide, "GitHub项目地址 + 个人联系方式\n完整代码与分析报告见GitHub仓库", 2, 5, 6, 1, size=16, color=RGBColor(100,100,100))

prs.save("运营商用户行为分析项目汇报.pptx")
print("PPT已生成：运营商用户行为分析项目汇报.pptx")